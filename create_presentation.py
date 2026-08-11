import os
import sys
import subprocess

def install_pptx():
    try:
        import pptx
    except ImportError:
        print("Installing python-pptx...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])

install_pptx()

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Define colors
    tecnm_blue = RGBColor(27, 57, 106)
    tecnm_gray = RGBColor(128, 128, 128)
    tecnm_gold = RGBColor(191, 151, 80)
    
    # 1. Title Slide
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Sistema Integral de Administración (SIA)"
    title.text_frame.paragraphs[0].font.color.rgb = tecnm_blue
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Definición de Alcance del Proyecto y Módulos Activos\nTecNM Campus Los Cabos"
    subtitle.text_frame.paragraphs[0].font.color.rgb = tecnm_gray
    
    # 2. Objetivo
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Objetivo del Proyecto"
    title.text_frame.paragraphs[0].font.color.rgb = tecnm_blue
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Modernizar y centralizar los servicios clave del instituto para mejorar la atención al estudiante y la toma de decisiones directiva."
    p = tf.add_paragraph()
    p.text = "Enfoque principal en:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Salud e Integridad (Servicios Médicos y Lactario)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Apoyo Académico (Biblioteca)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Gestión de la Calidad (Encuestas, Quejas, Avisos)"
    p.level = 2
    
    # 3. Modulos Principales
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Módulos que se MANTIENEN (Core & Operación)"
    title.text_frame.paragraphs[0].font.color.rgb = tecnm_blue
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Servicios Estudiantiles Críticos:"
    
    p = tf.add_paragraph()
    p.text = "Biblioteca:"
    p.font.bold = True
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Gestión de inventario, préstamos digitales, control de multas y visitas."
    p2.level = 2
    
    p = tf.add_paragraph()
    p.text = "Servicios Médicos:"
    p.font.bold = True
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Agendado de citas (Médico y Psicológico), expedientes clínicos seguros y control de turnos."
    p2.level = 2
    
    p = tf.add_paragraph()
    p.text = "Lactario:"
    p.font.bold = True
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Gestión de reservas de espacios con vinculación a apoyo médico en caso de requerirse."
    p2.level = 2

    # 4. Modulos Calidad
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Módulos que se MANTIENEN (Calidad & Directivo)"
    title.text_frame.paragraphs[0].font.color.rgb = tecnm_blue
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Evaluación Continua y Difusión:"
    
    p = tf.add_paragraph()
    p.text = "Encuestas Institucionales:"
    p.font.bold = True
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Segmentadas por rol (estudiantes, docentes, etc.), opcionales u obligatorias (bloqueo de accesos hasta responder)."
    p2.level = 2
    
    p = tf.add_paragraph()
    p.text = "Buzón de Quejas y Sugerencias:"
    p.font.bold = True
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Recepción de evidencias fotográficas, modo anónimo y gestión de tickets de atención."
    p2.level = 2
    
    p = tf.add_paragraph()
    p.text = "Dashboard Directivo (Reportes):"
    p.font.bold = True
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Métricas consolidadas de visitas, encuestas y poblacion en tiempo real para toma de decisiones."
    p2.level = 2

    # 5. Modulos Pausados
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Módulos EXCLUIDOS (Para fases futuras)"
    title.text_frame.paragraphs[0].font.color.rgb = tecnm_gray
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Los siguientes módulos no se incluirán en el lanzamiento inicial para enfocar esfuerzos en los servicios prioritarios:"
    
    p = tf.add_paragraph()
    p.text = "Cafetería (Sistema de menús y pedidos)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Comunidad (Foros de discusión de alumnos)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Aula (Gestión de tareas y clases)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Test Vocacional (Uso externo)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Proceso de Admisión (Aspirantes)"
    p.level = 1

    # 6. Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Impacto Inmediato"
    title.text_frame.paragraphs[0].font.color.rgb = tecnm_blue
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Beneficios del enfoque priorizado:"
    
    p = tf.add_paragraph()
    p.text = "Asegura la estabilidad y escalabilidad de los procesos que más valor aportan diariamente (Salud y Biblioteca)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Ofrece a Dirección General visibilidad inmediata del estado de los servicios (Módulo Calidad/Reportes)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Mantiene los costos de infraestructura optimizados."
    p.level = 1
    
    pptx_path = os.path.join(os.getcwd(), 'Presentacion_Alcances_SIA.pptx')
    prs.save(pptx_path)
    print(f"Presentation saved successfully to: {pptx_path}")

if __name__ == "__main__":
    create_presentation()
